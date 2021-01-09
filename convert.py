import os, sys

from PIL import Image
from pdf2image import convert_from_path
from docx2pdf import convert
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

file_to_convert = sys.argv[1]

print()
print("Converting file: " + file_to_convert)
print()

# Check if file is docx, then we need to convert it to pdf first
if file_to_convert.endswith('.docx'):
	print("Converting Word document to pdf first ...")
	convert(file_to_convert)
	base, _ = os.path.splitext(file_to_convert)
	pdf_file = base + "." + "pdf"
else:
	pdf_file = file_to_convert

# Prep presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# Create working folder
base_name, _ = os.path.splitext(pdf_file)

# Convert PDF to list of images
print("Starting conversion...")
slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2)
print("...complete.")
print()

# Loop over slides
for i, slideimg in enumerate(slideimgs):
	if i % 10 == 0:
		print("Saving slide: " + str(i))

	imagefile = BytesIO()
	slideimg.save(imagefile, format='tiff')
	imagedata = imagefile.getvalue()
	imagefile.seek(0)
	width, height = slideimg.size

	# Set slide dimensions
	prs.slide_height = height * 9525
	prs.slide_width = width * 9525

	# Add slide
	slide = prs.slides.add_slide(blank_slide_layout)
	pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

# Save Powerpoint
print()
print("Saving file: " + base_name + ".pptx")
prs.save(base_name + '.pptx')
print("Conversion complete. :)")
print()
